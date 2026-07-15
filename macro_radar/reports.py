from __future__ import annotations

from io import BytesIO
from pathlib import Path

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from .analytics import build_summary, latest_per_indicator, prepare_frame, previous_change
from .models import INDICATORS


def build_excel(records: list[dict[str, object]]) -> bytes:
    frame = prepare_frame(records)
    export_frame = frame.copy()
    for column in ("fetched_at", "published_at", "created_at", "updated_at"):
        if column not in export_frame.columns:
            continue
        values = pd.to_datetime(export_frame[column], errors="coerce", utc=True)
        export_frame[column] = values.dt.tz_localize(None)
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        export_frame.to_excel(writer, sheet_name="All_Data", index=False)
        if not export_frame.empty:
            for code, group in export_frame.groupby("indicator_code"):
                group.to_excel(writer, sheet_name=str(code)[:31], index=False)
    return output.getvalue()


def _chart_image(frame: pd.DataFrame, codes: list[str], title: str) -> BytesIO:
    figure, axis = plt.subplots(figsize=(8.8, 4.4))
    for code in codes:
        series = frame.loc[frame["indicator_code"] == code].sort_values("period")
        if series.empty:
            continue
        axis.plot(
            series["period"],
            series["value"],
            marker="o",
            linewidth=2,
            label=INDICATORS.get(code, {}).get("label", code),
        )
    axis.set_title(title)
    axis.grid(alpha=0.25)
    axis.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
    figure.autofmt_xdate()
    if len(codes) > 1:
        axis.legend()
    figure.tight_layout()
    output = BytesIO()
    figure.savefig(output, format="png", dpi=160, bbox_inches="tight")
    plt.close(figure)
    output.seek(0)
    return output


def build_powerpoint(
    records: list[dict[str, object]],
    template_path: Path | None = None,
) -> bytes:
    frame = prepare_frame(records)
    if frame.empty:
        raise ValueError("Cannot create a report without observations")

    use_template = template_path and template_path.exists() and template_path.stat().st_size > 0
    presentation = Presentation(str(template_path)) if use_template else Presentation()
    latest = latest_per_indicator(frame)
    report_date = frame["period"].max().date().isoformat()

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Macro Radar Kazakhstan"
    slide.placeholders[1].text = f"Данные на {report_date}"

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Macro Dashboard"
    codes = [
        code
        for code in ("USD_KZT", "CPI_YOY", "NBK_BASE_RATE", "TONIA", "BRENT_USD")
        if code in latest.index
    ]
    table = slide.shapes.add_table(
        rows=len(codes) + 1,
        cols=5,
        left=Inches(0.45),
        top=Inches(1.35),
        width=Inches(9.0),
        height=Inches(0.55 * (len(codes) + 1)),
    ).table
    for index, title in enumerate(("Показатель", "Значение", "Изменение", "Период", "Источник")):
        table.cell(0, index).text = title
    for row_index, code in enumerate(codes, start=1):
        row = latest.loc[code]
        change = previous_change(frame, code)
        table.cell(row_index, 0).text = INDICATORS[code]["label"]
        table.cell(row_index, 1).text = f"{float(row['value']):.2f} {row['unit']}"
        table.cell(row_index, 2).text = "—" if change is None else f"{change:+.2f}"
        table.cell(row_index, 3).text = pd.Timestamp(row["period"]).date().isoformat()
        table.cell(row_index, 4).text = str(row["source"])

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = "\n".join(f"• {line}" for line in build_summary(frame))

    chart_specs = [
        (["USD_KZT"], "USD/KZT"),
        (
            ["CPI_YOY", "CPI_FOOD_YOY", "CPI_NON_FOOD_YOY", "CPI_SERVICES_YOY"],
            "Инфляция и компоненты",
        ),
        (["BRENT_USD"], "Brent"),
    ]
    for chart_codes, title in chart_specs:
        if not any(code in latest.index for code in chart_codes):
            continue
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        image = _chart_image(frame, chart_codes, title)
        slide.shapes.add_picture(image, Inches(0.8), Inches(1.35), width=Inches(8.4))

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()
